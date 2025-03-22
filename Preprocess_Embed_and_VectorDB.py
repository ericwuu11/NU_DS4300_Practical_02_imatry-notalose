import os
import re
import time
import psutil
import numpy as np
import pandas as pd
import redis
import chromadb
import faiss
from sentence_transformers import SentenceTransformer
from pdfminer.high_level import extract_text
from pptx import Presentation
from sklearn.metrics.pairwise import cosine_similarity

# Defining the directory paths
corpus_dir = "C:/Users/jenny/Downloads/DS4300_25/Practical02/Corpus"
chunked_text_file = os.path.join(corpus_dir, "chunked_text_results.csv")
embedding_output_file = os.path.join(corpus_dir, "embedding_model_comparison.csv")
vector_db_output_file = os.path.join(corpus_dir, "vector_db_comparison.csv")

# List of selected input files within the "Corpus" folder
file_paths = [
    os.path.join(corpus_dir, "02 - Foundations.pptx"),
    os.path.join(corpus_dir, "03 - Moving Beyond the Relational Model.pptx"),
    os.path.join(corpus_dir, "04 - Data Replication.pptx"),
    os.path.join(corpus_dir, "05 - NoSQL Intro + KV DBs.pptx"),
    os.path.join(corpus_dir, "ICS 46 Spring 2022, Notes and Examples_ AVL Trees.pdf")
]

# Function to measure current memory usage in MB
def get_memory_usage():
    return psutil.Process().memory_info().rss / (1024 * 1024)

# Function to clean text by converting to lowercase, removing extra spaces, and removing punctuation
def clean_text(text):
    text = text.lower()
    text = re.sub(r'\s+', ' ', text)
    text = re.sub(r'[^\w\s]', '', text)
    return text.strip()

# Function to split text into overlapping chunks of specified sizes
def chunk_text(text, chunk_size, overlap_size):
    tokens = text.split()
    chunks = []
    for i in range(0, len(tokens), chunk_size - overlap_size):
        chunk = tokens[i:i + chunk_size]
        chunks.append(" ".join(chunk))
    return chunks

# Function to extract text from multiple files
def extract_text_from_files(file_paths):
    all_text = ""

    # Processing each file based on its format
    for path in file_paths:
        if path.endswith(".pptx"):
            prs = Presentation(path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        all_text += shape.text + " "
        elif path.endswith(".pdf"):
            all_text += extract_text(path) + " "

    return clean_text(all_text)

# Function to extract, clean, chunk, and save text to a CSV file
def process_and_chunk_text():
    cleaned_text = extract_text_from_files(file_paths)

    # Defining different chunk sizes and overlaps for experimentation
    chunk_sizes = [200, 500, 1000]
    overlap_sizes = [0, 50, 100]

    chunk_results = {}
    for chunk_size in chunk_sizes:
        for overlap in overlap_sizes:
            key = f"Chunk {chunk_size} | Overlap {overlap}"
            chunk_results[key] = chunk_text(cleaned_text, chunk_size, overlap)

    # Saving chunked text results to a CSV file
    df_chunks = pd.DataFrame.from_dict(chunk_results, orient="index").transpose()
    df_chunks.to_csv(chunked_text_file, index=False)

# Function to generate and compare embeddings using different models
def generate_embeddings():
    # Handling missing values
    df_chunks = pd.read_csv(chunked_text_file).fillna("")
    # Converting text to list format
    chunks = df_chunks.iloc[:, 0].astype(str).tolist()

    # Defining embedding models to test
    models = {
        "MiniLM": "sentence-transformers/all-MiniLM-L6-v2",
        "MPNet": "sentence-transformers/all-mpnet-base-v2",
        "GTE-Small": "thenlper/gte-small"
    }

    embedding_results = {}

    # Iterating over each model to generate embeddings
    for model_name, model_path in models.items():
        model = SentenceTransformer(model_path)

        start_time = time.time()
        start_memory = get_memory_usage()

        try:
            embeddings = model.encode(chunks, convert_to_numpy=True, batch_size=4)
        except Exception as e:
            print(f"Error while encoding with {model_name}: {e}")
            continue

        end_time = time.time()
        end_memory = get_memory_usage()

        # Computing similarity between embeddings
        similarity_matrix = cosine_similarity(embeddings)
        avg_similarity = np.mean(similarity_matrix)

        # Storing embedding results
        embedding_results[model_name] = {
            "Embedding Shape": embeddings.shape,
            "Computation Time (s)": round(end_time - start_time, 2),
            "Memory Usage (MB)": round(end_memory - start_memory, 2),
            "Average Cosine Similarity": round(avg_similarity, 4),
            "Retrieval Quality (Qualitative)": "Good for fast retrieval" if model_name == "MiniLM" else "Better for deeper understanding"
        }

    # Saving embedding comparison results to a CSV file
    df_results = pd.DataFrame.from_dict(embedding_results, orient="index")
    df_results.to_csv(embedding_output_file)

# Function to compare different vector databases (Redis, ChromaDB, FAISS)
def compare_vector_databases():
    model = SentenceTransformer("sentence-transformers/all-MiniLM-L6-v2")
    num_vectors = 1000
    embedding_dim = 384
    sample_texts = [f"Sample text {i}" for i in range(num_vectors)]
    embeddings = model.encode(sample_texts, convert_to_numpy=True).astype("float32")

    results = {}

    # Storing embeddings in Redis
    redis_client = redis.Redis(host="localhost", port=6379, db=0)
    redis_pipeline = redis_client.pipeline()

    start_mem = get_memory_usage()
    start_time = time.time()
    for i, vector in enumerate(embeddings):
        redis_pipeline.set(f"vector:{i}", vector.tobytes())
    redis_pipeline.execute()
    indexing_time = time.time() - start_time
    end_mem = get_memory_usage()

    start_time = time.time()
    redis_client.get("vector:0")
    query_time = time.time() - start_time

    results["Redis"] = {
        "Indexing Time (s)": round(indexing_time, 3),
        "Query Time (s)": round(query_time, 3),
        "Memory Usage (MB)": round(end_mem - start_mem, 3),
    }

    # Storing embeddings in ChromaDB
    chroma_client = chromadb.PersistentClient(path="./chroma_db")
    collection = chroma_client.get_or_create_collection(name="vector_test")

    start_mem = get_memory_usage()
    start_time = time.time()
    for i, vector in enumerate(embeddings):
        collection.add(ids=[str(i)], embeddings=[vector.tolist()])
    indexing_time = time.time() - start_time
    end_mem = get_memory_usage()

    start_time = time.time()
    collection.query(query_embeddings=[embeddings[0].tolist()], n_results=1)
    query_time = time.time() - start_time

    results["Chroma"] = {
        "Indexing Time (s)": round(indexing_time, 3),
        "Query Time (s)": round(query_time, 3),
        "Memory Usage (MB)": round(end_mem - start_mem, 3),
    }

    # Storing embeddings in FAISS
    start_mem = get_memory_usage()
    start_time = time.time()
    index = faiss.IndexFlatL2(embedding_dim)
    index.add(embeddings)
    indexing_time = time.time() - start_time
    end_mem = get_memory_usage()

    start_time = time.time()
    index.search(np.expand_dims(embeddings[0], axis=0), k=1)
    query_time = time.time() - start_time

    results["FAISS"] = {
        "Indexing Time (s)": round(indexing_time, 3),
        "Query Time (s)": round(query_time, 3),
        "Memory Usage (MB)": round(end_mem - start_mem, 3),
    }

    # Saving vector database comparison results to a CSV file
    df_results = pd.DataFrame.from_dict(results, orient="index")
    df_results.to_csv(vector_db_output_file)

# Executing all processing steps
if __name__ == "__main__":
    process_and_chunk_text()
    generate_embeddings()
    compare_vector_databases()
